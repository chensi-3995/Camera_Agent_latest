from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt


ROOT_DIR = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT_DIR / "docs" / "opening_defense_camera_agent_local.pptx"
TITLE = "基于多模态智能体的本地化视频安防监控系统设计与实现"
SUBTITLE = "开题答辩 | 研究院内网直连版本"
AUTHOR_LINE = "答辩人：陈思        导师：待补充"
DATE_LINE = "2026年4月"

SECTIONS = [
    ("研究背景", "Background"),
    ("系统方案", "Architecture"),
    ("智能体设计", "Agent"),
    ("阶段进展", "Progress"),
    ("后续计划", "Plans"),
]

GREEN = RGBColor(17, 122, 75)
LIGHT_GREEN = RGBColor(224, 243, 232)
DEEP_GREEN = RGBColor(0, 95, 62)
TEXT = RGBColor(40, 40, 40)
MUTED = RGBColor(110, 120, 120)
BG = RGBColor(248, 250, 248)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 226, 216)
BLUE = RGBColor(47, 111, 235)
ORANGE = RGBColor(219, 138, 0)
RED = RGBColor(196, 64, 64)


def find_template_path() -> Path:
    candidates = [
        path
        for path in (Path.home() / "Desktop").glob("*XRQ-2.pptx")
        if not path.name.startswith("~$")
    ]
    if not candidates:
        raise FileNotFoundError("未在桌面找到模板文件 *XRQ-2.pptx")
    return candidates[0]


def remove_all_slides(prs: Presentation) -> None:
    sld_id_list = prs.slides._sldIdLst
    while len(sld_id_list):
        rel_id = sld_id_list[0].rId
        prs.part.drop_rel(rel_id)
        del sld_id_list[0]


def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else shape_type,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.2)
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=20,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Microsoft YaHei",
    vertical_anchor=MSO_ANCHOR.MIDDLE,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = vertical_anchor
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(3)
    frame.margin_bottom = Pt(3)
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline_block(slide, left, top, width, height, lines, *, size=18, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Pt(10)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(8)
    frame.margin_bottom = Pt(6)
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = line
        para.level = 0
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_top_nav(slide, active_index: int, page_no: int, title: str, subtitle: str):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(2.2), WHITE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(2.18), Cm(33.87), Cm(0.05), GREEN)

    x = Cm(6.9)
    tab_widths = [Cm(4.4), Cm(4.4), Cm(4.4), Cm(4.4), Cm(4.8)]
    for idx, (cn, en) in enumerate(SECTIONS):
        fill = GREEN if idx == active_index else RGBColor(242, 246, 244)
        font_color = WHITE if idx == active_index else MUTED
        add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            x,
            Cm(0.38),
            tab_widths[idx],
            Cm(1.3),
            fill,
        )
        add_text(slide, x, Cm(0.43), tab_widths[idx], Cm(0.55), cn, size=15, color=font_color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x, Cm(0.98), tab_widths[idx], Cm(0.45), en, size=9, color=font_color, align=PP_ALIGN.CENTER)
        x += tab_widths[idx] + Cm(0.05)

    add_text(slide, Cm(0.8), Cm(0.25), Cm(5.5), Cm(0.75), "智能视频监控 Agent 系统", size=17, color=GREEN, bold=True)
    add_text(slide, Cm(0.8), Cm(1.0), Cm(5.5), Cm(0.45), "本地直连版开题答辩", size=10, color=MUTED)
    add_text(slide, Cm(0.8), Cm(2.45), Cm(14), Cm(0.8), title, size=22, color=TEXT, bold=True)
    add_text(slide, Cm(0.8), Cm(3.2), Cm(16), Cm(0.5), subtitle, size=10, color=GREEN)
    add_text(slide, Cm(30.8), Cm(17.2), Cm(1.8), Cm(0.5), str(page_no), size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body, *, fill=WHITE, title_color=GREEN):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line_color=LINE, radius=True)
    add_text(slide, left + Cm(0.35), top + Cm(0.25), width - Cm(0.7), Cm(0.7), title, size=16, color=title_color, bold=True)
    add_multiline_block(slide, left + Cm(0.25), top + Cm(1.0), width - Cm(0.5), height - Cm(1.2), body, size=14)


def add_metric_box(slide, left, top, width, height, value, label, color):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, WHITE, line_color=color, radius=True)
    add_text(slide, left, top + Cm(0.2), width, Cm(0.9), value, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Cm(0.15), top + Cm(1.15), width - Cm(0.3), Cm(0.6), label, size=12, color=TEXT, align=PP_ALIGN.CENTER)


def add_process_arrow(slide, x, y, w, h, title, body, fill_color):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h, fill_color)
    add_text(slide, x + Cm(0.2), y + Cm(0.15), w - Cm(0.7), Cm(0.7), title, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Cm(0.15), y + Cm(0.85), w - Cm(0.7), Cm(1.6), body, size=11, color=WHITE, align=PP_ALIGN.CENTER)


def add_slide_title(slide, title: str, subtitle: str, section_index: int, page_no: int):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_top_nav(slide, section_index, page_no, title, subtitle)


def extract_template_images(template: Presentation) -> dict[str, list[tuple[bytes, int, int, int, int]]]:
    image_data: dict[str, list[tuple[bytes, int, int, int, int]]] = {"cover": [], "thanks": []}
    for shape in template.slides[0].shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_data["cover"].append((shape.image.blob, shape.left, shape.top, shape.width, shape.height))
    for shape in template.slides[-1].shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_data["thanks"].append((shape.image.blob, shape.left, shape.top, shape.width, shape.height))
    return image_data


def build_cover(prs: Presentation, images: dict[str, list[tuple[bytes, int, int, int, int]]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 249, 246)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(1.3), GREEN)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(11.5), Cm(33.87), Cm(7.55), RGBColor(241, 247, 243))
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Cm(16.0), Cm(9.7), Cm(3.6), Cm(2.8), GREEN)
    add_text(slide, Cm(2.0), Cm(4.2), Cm(29), Cm(1.5), TITLE, size=28, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(4.0), Cm(6.0), Cm(25), Cm(0.8), SUBTITLE, size=18, color=GREEN, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(8.8), Cm(7.45), Cm(16.2), Cm(0.06), GREEN)
    add_text(slide, Cm(7.5), Cm(8.0), Cm(18.8), Cm(0.8), AUTHOR_LINE, size=16, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(7.5), Cm(8.9), Cm(18.8), Cm(0.7), DATE_LINE, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    for blob, left, top, width, height in images.get("cover", []):
        slide.shapes.add_picture(BytesIO(blob), left, top, width=width, height=height)


def build_agenda(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "答辩内容目录", "围绕项目背景、系统实现、智能体架构、阶段进展与后续计划展开", 0, page_no)

    cards = [
        ("01 研究背景与问题定义", ["传统视频全量上云成本高", "人工巡检难以处理长时间、多摄像头视频", "需要本地化、可交互、可追溯的智能安防方案"], GREEN),
        ("02 本地直连版系统方案", ["两路 RTSP 摄像头接入研究院内网", "本地笔记本完成录制、抽帧、分析与存储", "Web 端统一展示关键帧、问答、日报和告警"], BLUE),
        ("03 安防智能体设计", ["重点讲 LangGraph 多节点工作流", "感知-记忆-推理-执行闭环", "多模型分工协作与结构化输出约束"], ORANGE),
        ("04 当前阶段进展与验证", ["已完成两路摄像头、本地 UI、飞书、SQLite/Qdrant", "已接入 Qwen2.5-VL、Qwen2.5、Qwen3-Embedding", "评估当前效果与尚待优化问题"], RED),
        ("05 后续语音模块与知识库计划", ["大厅麦克风唤醒安防智能体", "构建人员特征本地知识库", "实现‘今天我有没有来过’类自然问答"], DEEP_GREEN),
    ]
    y_positions = [Cm(4.3), Cm(7.0), Cm(9.7), Cm(12.4), Cm(15.1)]
    for (title, lines, color), y in zip(cards, y_positions):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.2), y, Cm(31.0), Cm(2.1), WHITE, line_color=LINE, radius=True)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.5), y + Cm(0.32), Cm(5.8), Cm(1.45), color, radius=True)
        add_text(slide, Cm(1.5), y + Cm(0.38), Cm(5.8), Cm(1.1), title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Cm(7.8), y + Cm(0.3), Cm(23.5), Cm(1.5), "  ｜  ".join(lines), size=13, color=TEXT)


def build_background(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "研究背景与问题定义", "目标是在研究院内网环境下构建低成本、可检索、可问答的本地化视频安防智能体", 0, page_no)

    add_card(slide, Cm(0.9), Cm(4.4), Cm(10.4), Cm(6.0), "痛点 1：全量视频直接上云成本高", ["多路摄像头持续产生日志与视频流", "带宽、云端存储、逐帧推理算力成本高", "大量静态画面造成重复传输和重复计算"], fill=RGBColor(245, 250, 246))
    add_card(slide, Cm(11.7), Cm(4.4), Cm(10.4), Cm(6.0), "痛点 2：事后检索依赖人工回看", ["长时间录像只能按时间线手工拖拽筛查", "难以按‘黑衣人员/几人路过/高危事件’语义查询", "缺少面向安保值班人员的自然语言入口"], fill=RGBColor(244, 248, 255), title_color=BLUE)
    add_card(slide, Cm(22.5), Cm(4.4), Cm(10.4), Cm(6.0), "痛点 3：告警与知识沉淀不足", ["异常发现、告警推送、日报总结之间未形成闭环", "缺少可长期积累的本地人员/事件知识库", "难以把‘看见了什么’转为可复用的业务记忆"], fill=RGBColor(255, 249, 238), title_color=ORANGE)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(0.9), Cm(11.1), Cm(32.0), Cm(6.1), LIGHT_GREEN, line_color=GREEN, radius=True)
    add_text(slide, Cm(1.5), Cm(11.55), Cm(31.0), Cm(0.8), "拟解决的核心问题", size=20, color=GREEN, bold=True)
    add_multiline_block(
        slide,
        Cm(1.6),
        Cm(12.4),
        Cm(30.0),
        Cm(4.2),
        [
            "1. 如何在本地设备上对两路摄像头视频做低冗余关键帧提取，避免把大量无效静态画面直接交给大模型？",
            "2. 如何用多模态智能体将关键帧分析、结构化入库、语义检索、自然语言问答、飞书告警串成可维护的闭环工作流？",
            "3. 如何在后续接入语音模块与人员知识库后，让系统支持‘今天谁来过/我有没有出现过’这类更自然的大厅交互？",
        ],
        size=16,
    )


def build_local_architecture(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "本地直连版系统架构", "只讲研究院内网直连方案：摄像头与本地工作站同网段，视频采集、分析、存储和交互均在本地闭环", 1, page_no)

    stages = [
        ("Cam_01 / Cam_02", "TP-Link 网络摄像头\nRTSP 子码流\n研究院内网直连", GREEN),
        ("视频录制层", "按片段录制 MP4\n本地缓存原始视频\n支持 Web 实时预览", BLUE),
        ("关键帧提取层", "OpenCV 运动检测\n窗口内峰值帧选择\n重复帧过滤与事件合并", ORANGE),
        ("多模态分析层", "Qwen2.5-VL 输出结构化 JSON\nPydantic 字段校验\n风险等级与事件描述生成", RED),
        ("记忆与交互层", "SQLite + Qdrant + Embedding\n自然语言问答\n全天总结与飞书告警", DEEP_GREEN),
    ]
    x = Cm(0.9)
    for title, body, color in stages:
        add_process_arrow(slide, x, Cm(5.4), Cm(6.3), Cm(4.4), title, body, color)
        x += Cm(6.35)

    add_card(slide, Cm(1.0), Cm(11.1), Cm(15.4), Cm(5.5), "当前本地直连版的技术边界", ["本地笔记本与两路摄像头必须处于同一研究院内网。", "当前不强调毫秒级实时，而是以“录制片段 + 抽帧分析 + 准实时展示/告警”为主。", "系统后续扩展优先围绕“智能体工作流稳定性、事件知识库、语音入口”推进。"])
    add_card(slide, Cm(16.8), Cm(11.1), Cm(15.8), Cm(5.5), "四个已落地业务能力", ["关键帧分析报告：按日期、摄像头、时间轴查看关键帧与大模型描述。", "高危异常飞书报警：高风险事件触发后推送图片与文字。", "全天日志 AI 总结：按早晨/下午/晚上/凌晨四段生成日报。", "自然语言监控对话：基于 SQLite + Qdrant 检索历史事件并用本地 Qwen2.5 组织回答。"])


def build_agent_loop(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "安防智能体：感知-记忆-推理-执行闭环", "把系统重点从“脚本顺序执行”升级为“可追踪状态流转 + 多模型协同 + 可扩展工具调用”的智能体架构", 2, page_no)

    loop_nodes = [
        (Cm(2.0), Cm(6.0), Cm(6.4), Cm(3.6), "感知模块", "读取 RTSP/视频片段\nOpenCV 抽取候选关键帧\n过滤静态冗余画面", GREEN),
        (Cm(9.7), Cm(6.0), Cm(6.4), Cm(3.6), "记忆模块", "SQLite 保存结构化事件\nQdrant 存语义向量\n为 RAG 与日报提供上下文", BLUE),
        (Cm(17.4), Cm(6.0), Cm(6.4), Cm(3.6), "推理模块", "Qwen2.5-VL 解析画面语义\nQwen2.5 生成问答/总结\nPydantic 约束输出格式", ORANGE),
        (Cm(25.1), Cm(6.0), Cm(6.4), Cm(3.6), "执行模块", "高危事件飞书告警\n普通事件入库归档\n日报生成并回写记忆库", RED),
    ]
    for left, top, width, height, title, body, color in loop_nodes:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, color, radius=True)
        add_text(slide, left + Cm(0.2), top + Cm(0.35), width - Cm(0.4), Cm(0.8), title, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, left + Cm(0.2), top + Cm(1.35), width - Cm(0.4), Cm(1.8), body, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, Cm(8.45), Cm(7.1), Cm(1.1), Cm(1.4), LIGHT_GREEN)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, Cm(16.15), Cm(7.1), Cm(1.1), Cm(1.4), LIGHT_GREEN)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, Cm(23.85), Cm(7.1), Cm(1.1), Cm(1.4), LIGHT_GREEN)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(2.0), Cm(11.2), Cm(29.5), Cm(5.0), WHITE, line_color=LINE, radius=True)
    add_text(slide, Cm(2.6), Cm(11.6), Cm(28.0), Cm(0.8), "为什么这一层要突出“智能体”而不是普通脚本？", size=18, color=GREEN, bold=True)
    add_multiline_block(
        slide,
        Cm(2.6),
        Cm(12.6),
        Cm(28.0),
        Cm(3.0),
        [
            "状态可追踪：每一帧/每一条用户问题都以 State 在节点间传递，方便调试、回溯和扩展。",
            "分支可控制：高危事件走“告警分支”，普通事件走“入库分支”，模型失败走“模型未加载成功”降级分支。",
            "记忆可累积：事件写入 SQLite/Qdrant 后，后续问答、日报、语音 Q&A 都能复用同一份本地知识资产。",
        ],
        size=15,
    )


def build_langgraph_workflow(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "LangGraph 多模型工作流设计", "多个模型分别负责图像理解、文本生成、向量检索，LangGraph 负责节点路由、状态传递和条件分支", 2, page_no)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.0), Cm(4.7), Cm(15.7), Cm(11.2), WHITE, line_color=LINE, radius=True)
    add_text(slide, Cm(1.6), Cm(5.0), Cm(14.5), Cm(0.7), "工作流 A：关键帧分析图", size=18, color=GREEN, bold=True)
    y = Cm(6.0)
    for title, color in [
        ("入队关键帧", GREEN),
        ("构建中文视觉提示词", BLUE),
        ("调用 Qwen2.5-VL", ORANGE),
        ("解析/校验 JSON", RED),
        ("事件合并 + 风险判断 + 写库/告警", DEEP_GREEN),
    ]:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(2.0), y, Cm(13.4), Cm(1.45), color, radius=True)
        add_text(slide, Cm(2.2), y + Cm(0.15), Cm(13.0), Cm(1.0), title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        y += Cm(1.75)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(17.2), Cm(4.7), Cm(15.7), Cm(11.2), WHITE, line_color=LINE, radius=True)
    add_text(slide, Cm(17.8), Cm(5.0), Cm(14.5), Cm(0.7), "工作流 B：自然语言监控对话图", size=18, color=GREEN, bold=True)
    y = Cm(6.0)
    for title, color in [
        ("用户问题 + 历史会话", GREEN),
        ("问题改写/时间与摄像头约束解析", BLUE),
        ("SQLite 精确筛选 + Qdrant 语义排序", ORANGE),
        ("事件匹配/统计聚合", RED),
        ("Qwen2.5 生成可读回答 + 返回关键帧", DEEP_GREEN),
    ]:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(18.2), y, Cm(13.4), Cm(1.45), color, radius=True)
        add_text(slide, Cm(18.4), y + Cm(0.15), Cm(13.0), Cm(1.0), title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        y += Cm(1.75)


def build_memory_rag(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "记忆与 RAG：本地结构化库 + 向量库双存储", "让“看过的视频”沉淀为可查询、可总结、可对话、可持续扩展的本地知识库", 2, page_no)

    add_card(slide, Cm(1.2), Cm(5.0), Cm(10.0), Cm(10.5), "SQLite 结构化记忆", ["tasks：一次分析任务的起止时间、触发类型、任务状态。", "camera_runs：每路摄像头录制出的片段路径、帧数、FPS、状态。", "events：关键帧事件时间、风险等级、人数、动作、衣着颜色、图片/视频路径。", "summaries：按日期保存全天总结及飞书推送状态。"])
    add_card(slide, Cm(12.0), Cm(5.0), Cm(10.0), Cm(10.5), "Qdrant 语义向量记忆", ["把事件描述、风险类型、摄像头、时间上下文编码为向量。", "支持“黑衣服的人”“有人徘徊吗”这类语义模糊查询。", "与 SQLite 精确过滤组合，实现“先硬约束、再语义排序”。", "后续人员知识库也可沿用同一套向量检索接口。"], fill=RGBColor(245, 248, 255), title_color=BLUE)
    add_card(slide, Cm(22.8), Cm(5.0), Cm(9.8), Cm(10.5), "Qwen3-Embedding + Qwen2.5", ["Qwen3-Embedding：本地部署，负责日志向量化。", "Qwen2.5：读取检索结果后，按安防汇报口吻生成回答/日报。", "Prompt 约束：只基于检索日志回答；无记录则返回“未在监控记录中发现相关异常”。"], fill=RGBColor(255, 249, 238), title_color=ORANGE)


def build_progress(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "目前项目进度", "当前已形成“本地直连摄像头采集 + 智能体分析 + 本地记忆与问答 + Web 可视化”的端到端原型", 3, page_no)

    headers = ["模块", "当前进度", "已实现内容", "下一步优化"]
    rows = [
        ["摄像头接入与录制", "已完成", "Cam_01 / Cam_02 RTSP 接入、片段录制、网页预览", "提升掉线重连与长时间运行稳定性"],
        ["关键帧提取", "已完成原型", "运动检测、窗口峰值帧、重复帧过滤、事件合并", "进一步减少无人物重复帧，提高人员经过召回率"],
        ["关键帧分析智能体", "已接入", "Qwen2.5-VL + 中文结构化 JSON + 风险分级", "优化提示词和字段质量，增强人物动作/衣着描述稳定性"],
        ["自然语言对话", "已上线", "基于 SQLite + Qdrant 的历史事件检索与流式回答", "增强复杂时间表达、统计类问答、多人/人物属性推理"],
        ["全天日志总结", "已上线", "按四时段总结 + 全天综合总结 + 飞书推送", "强化重点异常提炼和巡检建议生成"],
        ["前端交互界面", "已上线", "四个 Tag 页面、聊天历史、流式输出、关键帧时间轴", "继续修正聊天气泡细节，补充语音入口和人员库页面"],
    ]
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Cm(0.9), Cm(4.9), Cm(32.0), Cm(11.0)).table
    widths = [Cm(5.2), Cm(3.2), Cm(13.5), Cm(10.1)]
    for idx, width in enumerate(widths):
        table.columns[idx].width = width
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.runs[0].font.name = "Microsoft YaHei"
        para.runs[0].font.size = Pt(14)
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = WHITE
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 250, 246) if row_idx % 2 else WHITE
            cell.text = value
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if col_idx < 2 else PP_ALIGN.LEFT
            para.runs[0].font.name = "Microsoft YaHei"
            para.runs[0].font.size = Pt(12)
            para.runs[0].font.color.rgb = GREEN if col_idx == 1 and value.startswith("已") else TEXT


def build_validation(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "初步验证结果与当前问题", "以功能闭环验证为主，已能完成事件入库、语义查询、日报生成和飞书告警，但仍需持续提升识别质量与工程稳定性", 3, page_no)

    add_metric_box(slide, Cm(1.1), Cm(4.8), Cm(6.0), Cm(2.3), "2 路", "内网 RTSP 摄像头接入", GREEN)
    add_metric_box(slide, Cm(7.6), Cm(4.8), Cm(6.0), Cm(2.3), "4 类", "已落地核心业务功能", BLUE)
    add_metric_box(slide, Cm(14.1), Cm(4.8), Cm(6.0), Cm(2.3), "3 模型", "VL / 文本 / Embedding 协同", ORANGE)
    add_metric_box(slide, Cm(20.6), Cm(4.8), Cm(6.0), Cm(2.3), "4 表", "SQLite 事件/任务/录制/日报", RED)
    add_metric_box(slide, Cm(27.1), Cm(4.8), Cm(5.7), Cm(2.3), "1 套", "本地 Web 操作台", DEEP_GREEN)
    add_card(slide, Cm(1.0), Cm(7.8), Cm(15.5), Cm(8.0), "已验证的闭环能力", ["网页端可启动任务、查看两路画面、按日期查看关键帧时间轴与分析描述。", "高危事件可触发飞书推送，普通事件结构化入 SQLite 并写入向量库。", "自然语言对话支持流式输出、历史会话记录和关键帧图片引用。", "日报模块可按早晨/下午/晚上/凌晨四个时段总结并生成综合研判文本。"])
    add_card(slide, Cm(17.0), Cm(7.8), Cm(15.8), Cm(8.0), "当前暴露的主要问题", ["关键帧提取仍可能出现“无人重复帧偏多”或“某路人员经过漏提取”的情况。", "跨帧同一事件合并还需要继续调时间阈值、相似度阈值和代表帧选择策略。", "统计型问答虽然已做 SQL 优先，但人物属性字段质量仍依赖 VL 输出稳定性。", "下一阶段需要把“人是谁/是否是某人再次出现”从事件记忆进一步扩展到人员知识库记忆。"], title_color=RED)


def build_voice_module(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "后续扩展方向一：大厅语音唤醒与常见 Q&A", "在大厅放置高质量麦克风，通过语音唤醒接入安防智能体，实现面向值班/访客的自然语言问答入口", 4, page_no)

    x = [Cm(1.0), Cm(9.8), Cm(18.6), Cm(27.4)]
    w = [Cm(7.0), Cm(7.0), Cm(7.0), Cm(5.0)]
    cards = [
        ("大厅麦克风阵列", "唤醒词检测\n近场语音采集\n降噪/回声抑制", GREEN),
        ("ASR + 意图识别", "语音转文字\n时间/摄像头/人物条件抽取\n问答/查询/播报路由", BLUE),
        ("智能体检索与推理", "访问事件库/人员库\nRAG 检索\nQwen2.5 生成回答", ORANGE),
        ("语音播报", "TTS 回复\n前端同步显示\n必要时附关键帧", RED),
    ]
    for idx, ((title, body, color), left, width) in enumerate(zip(cards, x, w)):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Cm(5.0), width, Cm(3.5), color, radius=True)
        add_text(slide, left + Cm(0.2), Cm(5.4), width - Cm(0.4), Cm(0.8), title, size=19 if idx < 3 else 18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, left + Cm(0.2), Cm(6.3), width - Cm(0.4), Cm(1.5), body, size=13, color=WHITE, align=PP_ALIGN.CENTER)
        if idx < 3:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, left + width + Cm(0.3), Cm(6.0), Cm(1.2), Cm(1.5), LIGHT_GREEN)

    add_card(slide, Cm(1.0), Cm(9.6), Cm(15.5), Cm(5.7), "目标交互样例", ["“今天大厅有没有谁来过？” → 返回时间段、出现次数、对应摄像头和关键帧。", "“今天我有没有来过？” → 结合人员知识库中的个人描述/特征，检索相似事件后给出结果。", "“今天有没有陌生人多次徘徊？” → 从事件库中按动作类型、人数和时间聚合回答。"], fill=RGBColor(245, 250, 246))
    add_card(slide, Cm(17.0), Cm(9.6), Cm(15.4), Cm(5.7), "工程注意点", ["麦克风建议优先选带远场拾音、降噪和 USB/网口稳定接入能力的型号。", "语音模块只做入口层，后端仍复用现有 LangGraph 问答工作流和本地知识库。", "需增加唤醒词误触发控制、权限边界、日志留存和隐私合规说明。"], title_color=RED)


def build_person_kb(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "后续扩展方向二：构建人员本地知识库", "把“事件日志”进一步升级为“人员档案 + 历史出现轨迹 + 视觉/文本特征索引”的可问答本地知识库", 4, page_no)

    add_card(slide, Cm(1.0), Cm(5.0), Cm(10.2), Cm(10.5), "人员档案表", ["person_id：人员唯一编号", "name / nickname：姓名或备注名", "role：员工/访客/安保/未知", "appearance_text：身高、发型、衣着偏好、携带物等自然语言描述", "face/body embedding：可选视觉向量特征", "authorized_zone：允许出现区域"])
    add_card(slide, Cm(11.8), Cm(5.0), Cm(10.2), Cm(10.5), "事件-人员关联表", ["event_id ↔ person_id：把历史关键帧事件与某个已知人员候选关联起来", "match_score：视觉/文本相似度或规则匹配分数", "is_manual_verified：是否经过人工确认", "first_seen_at / last_seen_at：第一次/最近一次出现时间", "trajectory_summary：当天或一周内的出现轨迹摘要"], fill=RGBColor(245, 248, 255), title_color=BLUE)
    add_card(slide, Cm(22.6), Cm(5.0), Cm(10.2), Cm(10.5), "面向智能体的查询能力", ["“今天我有没有来过？”", "“张三今天最后一次出现在哪个摄像头？”", "“这周有没有陌生人连续三天晚上出现？”", "“穿黑色外套的人是不是上周二也出现过？”", "本质上是“摄像头事件记忆 + 人员知识库 + 语义检索 + 智能体推理”的组合。"], fill=RGBColor(255, 249, 238), title_color=ORANGE)


def build_schedule(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "后续研究计划与里程碑", "围绕“识别更准、记忆更强、交互更自然、工程更稳定”推进后续工作", 4, page_no)

    milestones = [
        ("4月", "优化关键帧提取与事件合并", ["复盘无人重复帧问题", "调参运动阈值/时间窗/相似度", "提升两路摄像头人员经过召回率"], GREEN),
        ("5月", "补强事件字段与统计问答", ["稳定输出人数/衣着/动作字段", "完善统计型 SQL 查询", "增强“上周几/某日晚上”时间解析"], BLUE),
        ("6月", "接入语音唤醒与大厅 Q&A", ["选型麦克风与唤醒链路", "ASR → 智能体 → TTS 闭环", "前端新增语音状态展示"], ORANGE),
        ("7月", "建设人员本地知识库", ["设计人员档案/事件关联表", "导入常见人员描述和特征", "支持“我有没有来过/某人轨迹”查询"], RED),
        ("8月", "系统评测与论文整理", ["长时间稳定性测试", "准确率/召回率/延迟指标评估", "整理开题后实验结果与论文结构"], DEEP_GREEN),
    ]
    x = Cm(1.0)
    for month, title, lines, color in milestones:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Cm(5.0), Cm(6.1), Cm(10.5), WHITE, line_color=LINE, radius=True)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Cm(0.4), Cm(5.4), Cm(5.3), Cm(1.4), color, radius=True)
        add_text(slide, x + Cm(0.5), Cm(5.55), Cm(1.4), Cm(0.9), month, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Cm(1.8), Cm(5.55), Cm(3.8), Cm(0.9), title, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_multiline_block(slide, x + Cm(0.4), Cm(7.3), Cm(5.3), Cm(7.5), lines, size=13)
        x += Cm(6.35)


def build_summary(prs: Presentation, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "总结与预期成果", "以“本地化安防智能体 + 可检索记忆库 + 自然语言/语音交互”为主线，形成可演示、可扩展、可论文落地的系统原型", 4, page_no)

    add_card(slide, Cm(1.0), Cm(5.2), Cm(10.0), Cm(9.8), "已完成基础", ["本地直连两路摄像头。", "关键帧提取、VL 分析、结构化入库。", "自然语言监控问答、四时段日报总结、飞书告警。", "SQLite + Qdrant 的本地记忆架构已经打通。"], fill=RGBColor(245, 250, 246))
    add_card(slide, Cm(11.9), Cm(5.2), Cm(10.0), Cm(9.8), "开题后重点突破", ["把 LangGraph 智能体工作流从“可跑”推进到“可解释、可评估、可长期稳定运行”。", "把事件库扩展为人员知识库，支持更贴近真实值班场景的自然问答。", "补充语音唤醒模块，让大厅场景可以直接对话使用系统。"], fill=RGBColor(244, 248, 255), title_color=BLUE)
    add_card(slide, Cm(22.8), Cm(5.2), Cm(10.0), Cm(9.8), "预期成果形式", ["一套面向研究院场景的本地化视频安防智能体原型系统。", "一套可复用的“摄像头事件记忆 + RAG + 语音入口”工程实现方案。", "一组围绕关键帧提取、事件问答准确性、告警时效性的实验结果。", "最终论文与答辩演示材料。"], fill=RGBColor(255, 249, 238), title_color=ORANGE)


def build_thanks(prs: Presentation, images: dict[str, list[tuple[bytes, int, int, int, int]]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 249, 246)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(1.25), GREEN)
    add_text(slide, Cm(4.0), Cm(6.8), Cm(25.8), Cm(1.5), "谢谢各位老师！", size=34, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Cm(4.0), Cm(8.5), Cm(25.8), Cm(0.9), "敬请批评指正", size=22, color=TEXT, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(11.0), Cm(10.0), Cm(11.8), Cm(0.06), GREEN)
    add_text(slide, Cm(4.0), Cm(10.6), Cm(25.8), Cm(0.8), AUTHOR_LINE, size=16, color=MUTED, align=PP_ALIGN.CENTER)
    for blob, left, top, width, height in images.get("thanks", []):
        slide.shapes.add_picture(BytesIO(blob), left, top, width=width, height=height)


def build_presentation() -> None:
    template_path = find_template_path()
    prs = Presentation(str(template_path))
    template_images = extract_template_images(prs)
    remove_all_slides(prs)

    build_cover(prs, template_images)
    build_agenda(prs, 2)
    build_background(prs, 3)
    build_local_architecture(prs, 4)
    build_agent_loop(prs, 5)
    build_langgraph_workflow(prs, 6)
    build_memory_rag(prs, 7)
    build_progress(prs, 8)
    build_validation(prs, 9)
    build_voice_module(prs, 10)
    build_person_kb(prs, 11)
    build_schedule(prs, 12)
    build_summary(prs, 13)
    build_thanks(prs, template_images)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
